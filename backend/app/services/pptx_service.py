import io
import logging
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.schemas.slide_schema import PresentationPayload, CoverData, BulletsData, SplitData, LayoutType

logger = logging.getLogger(__name__)

# 主题颜色配置
THEME_COLORS = {
    "default": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x1F, 0x29, 0x37),
        "text": RGBColor(0x37, 0x47, 0x51),
        "accent": RGBColor(0x3B, 0x82, 0xF6),
    },
    "dark": {
        "bg": RGBColor(0x1E, 0x29, 0x3B),
        "title": RGBColor(0xF1, 0xF5, 0xF9),
        "text": RGBColor(0xCB, 0xD5, 0xE1),
        "accent": RGBColor(0x60, 0xA5, 0xFA),
    },
    "corporate": {
        "bg": RGBColor(0xF8, 0xFA, 0xFC),
        "title": RGBColor(0x0F, 0x17, 0x2A),
        "text": RGBColor(0x1E, 0x40, 0xAF),
        "accent": RGBColor(0x1D, 0x4E, 0xD8),
    },
}


def _set_bg_color(slide, color: RGBColor):
    """设置幻灯片背景色"""
    from pptx.oxml.ns import qn
    from lxml import etree

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int, color: RGBColor, bold: bool = False,
                  align=PP_ALIGN.LEFT, font_name: str = "微软雅黑"):
    """添加文本框辅助函数"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text

    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.color.rgb = color
    font.bold = bold

    return txBox


def build_pptx(payload: PresentationPayload) -> io.BytesIO:
    """
    将 PresentationPayload 转换为 PPTX 二进制数据。
    返回 BytesIO 对象供路由层直接输出。
    """
    prs = Presentation()
    # 设置 16:9 比例（宽 33.87cm × 高 19.05cm）
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    colors = THEME_COLORS.get(payload.theme, THEME_COLORS["default"])

    for slide_payload in payload.slides:
        # 使用空白布局，手动布局所有元素
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 设置背景色
        _set_bg_color(slide, colors["bg"])

        if slide_payload.layout == LayoutType.COVER:
            _build_cover_slide(slide, slide_payload.data, colors, prs)
        elif slide_payload.layout == LayoutType.BULLETS:
            _build_bullets_slide(slide, slide_payload.data, colors, prs)
        elif slide_payload.layout == LayoutType.SPLIT:
            _build_split_slide(slide, slide_payload.data, colors, prs)
        else:
            logger.warning(f"未知的 layout 类型：{slide_payload.layout}，跳过")

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _build_cover_slide(slide, data: CoverData, colors: dict, prs: Presentation):
    """构建封面页"""
    w = prs.slide_width
    h = prs.slide_height

    # 顶部装饰条
    from pptx.util import Pt as PtUtil
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(0), Emu(int(h * 0.35)),
        w, Emu(int(h * 0.005))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["accent"]
    shape.line.fill.background()

    # 标题
    _add_text_box(
        slide, data.title,
        Emu(int(w * 0.1)), Emu(int(h * 0.25)),
        Emu(int(w * 0.8)), Emu(int(h * 0.2)),
        font_size=40, color=colors["title"],
        bold=True, align=PP_ALIGN.CENTER
    )

    # 副标题
    _add_text_box(
        slide, data.subtitle,
        Emu(int(w * 0.1)), Emu(int(h * 0.5)),
        Emu(int(w * 0.8)), Emu(int(h * 0.25)),
        font_size=20, color=colors["text"],
        bold=False, align=PP_ALIGN.CENTER
    )


def _build_bullets_slide(slide, data: BulletsData, colors: dict, prs: Presentation):
    """构建要点页"""
    w = prs.slide_width
    h = prs.slide_height

    # 顶部装饰条
    shape = slide.shapes.add_shape(
        1,
        Emu(0), Emu(0),
        w, Emu(int(h * 0.008))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["accent"]
    shape.line.fill.background()

    # 标题
    _add_text_box(
        slide, data.title,
        Emu(int(w * 0.07)), Emu(int(h * 0.06)),
        Emu(int(w * 0.86)), Emu(int(h * 0.15)),
        font_size=28, color=colors["title"],
        bold=True, align=PP_ALIGN.LEFT
    )

    # 分隔线
    line_shape = slide.shapes.add_shape(
        1,
        Emu(int(w * 0.07)), Emu(int(h * 0.22)),
        Emu(int(w * 0.86)), Emu(int(h * 0.003))
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = colors["accent"]
    line_shape.line.fill.background()

    # 要点列表
    bullet_height = Emu(int(h * 0.1))
    for i, bullet in enumerate(data.bullets):
        top = Emu(int(h * 0.27)) + i * bullet_height
        # 圆点装饰
        dot = slide.shapes.add_shape(
            9,  # 椭圆
            Emu(int(w * 0.07)), top + Emu(int(bullet_height * 0.3)),
            Emu(int(w * 0.012)), Emu(int(w * 0.012))
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors["accent"]
        dot.line.fill.background()

        _add_text_box(
            slide, bullet,
            Emu(int(w * 0.1)), top,
            Emu(int(w * 0.83)), bullet_height,
            font_size=16, color=colors["text"],
            bold=False, align=PP_ALIGN.LEFT
        )


def _build_split_slide(slide, data: SplitData, colors: dict, prs: Presentation):
    """构建分栏页"""
    w = prs.slide_width
    h = prs.slide_height

    # 顶部装饰条
    shape = slide.shapes.add_shape(
        1,
        Emu(0), Emu(0),
        w, Emu(int(h * 0.008))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["accent"]
    shape.line.fill.background()

    # 标题
    _add_text_box(
        slide, data.title,
        Emu(int(w * 0.07)), Emu(int(h * 0.06)),
        Emu(int(w * 0.86)), Emu(int(h * 0.15)),
        font_size=28, color=colors["title"],
        bold=True, align=PP_ALIGN.LEFT
    )

    # 左栏背景
    left_bg = slide.shapes.add_shape(
        1,
        Emu(int(w * 0.05)), Emu(int(h * 0.25)),
        Emu(int(w * 0.42)), Emu(int(h * 0.65))
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    left_bg.line.fill.background()

    # 右栏背景
    right_bg = slide.shapes.add_shape(
        1,
        Emu(int(w * 0.53)), Emu(int(h * 0.25)),
        Emu(int(w * 0.42)), Emu(int(h * 0.65))
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    right_bg.line.fill.background()

    # 左栏内容
    _add_text_box(
        slide, data.leftContent,
        Emu(int(w * 0.07)), Emu(int(h * 0.28)),
        Emu(int(w * 0.38)), Emu(int(h * 0.58)),
        font_size=14, color=colors["text"],
        bold=False, align=PP_ALIGN.LEFT
    )

    # 右栏内容
    _add_text_box(
        slide, data.rightContent,
        Emu(int(w * 0.55)), Emu(int(h * 0.28)),
        Emu(int(w * 0.38)), Emu(int(h * 0.58)),
        font_size=14, color=colors["text"],
        bold=False, align=PP_ALIGN.LEFT
    )
